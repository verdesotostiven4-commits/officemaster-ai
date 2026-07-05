from pathlib import Path
from pptx import Presentation


def create_powerpoint(prompt: str, title: str, style: str, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    file_path = output_dir / 'officemaster_powerpoint.pptx'

    presentation = Presentation()

    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = title or 'Presentacion generada por OfficeMaster AI'
    slide.placeholders[1].text = 'Estilo: ' + style

    sections = [
        ('Objetivo', 'Presentar el tema de forma clara, visual y profesional.'),
        ('Solicitud del usuario', prompt),
        ('Puntos principales', 'Aqui se colocaran las ideas clave generadas por IA.'),
        ('Desarrollo', 'Cada punto puede convertirse en una diapositiva resumida.'),
        ('Conclusion', 'Cierre claro con las ideas mas importantes.')
    ]

    for heading, body in sections:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = heading
        slide.placeholders[1].text = body

    presentation.save(file_path)
    return file_path
